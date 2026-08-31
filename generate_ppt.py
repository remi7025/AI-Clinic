"""Generate a human-style PowerPoint for the AI Clinic project."""

from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).parent
DATA_PATH = BASE / "data" / "compliance_dataset.json"
FIG_DIR = BASE / "latex" / "figures"
OUT = BASE / "AI_Clinic_Presentation.pptx"

NAVY = RGBColor(0x1D, 0x2A, 0x4D)
GOLD = RGBColor(0xC5, 0xA0, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
DARK = RGBColor(0x2A, 0x2A, 0x2A)
MED = RGBColor(0x4A, 0x4A, 0x4A)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)


def clean(text: str) -> str:
    return (
        text.replace("\u2018", "'")
        .replace("\u2019", "'")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2013", "-")
        .replace("\u2014", " - ")
    )


def bg(slide, color, left=0, top=0, width=None, height=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width or Inches(13.33), height or Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def box(slide, left, top, width, height, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tx


def bullets(slide, items, left, top, width, height, size=15, color=DARK, after=Pt(8)):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = clean(item)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = after
        p.level = 0
    return tx


def title_bar(slide, title: str):
    bg(slide, WHITE)
    bg(slide, NAVY, height=Inches(1.05))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), Inches(13.33), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    box(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.6), title, size=26, bold=True, color=WHITE)
    # footer
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), Inches(13.33), Inches(0.35))
    foot.fill.solid()
    foot.fill.fore_color.rgb = LIGHT
    foot.line.fill.background()
    box(
        slide,
        Inches(0.55),
        Inches(7.18),
        Inches(12),
        Inches(0.28),
        "AI Clinic  |  aivancity  |  AI Healthcare Compliance Across Countries",
        size=10,
        color=MED,
    )


def add_picture_safe(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        box(slide, left, top, width or Inches(6), height or Inches(3.5), f"[Missing figure: {path.name}]", size=12, color=RED)
        return
    if width and height:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(path), left, top, width=width)
    else:
        slide.shapes.add_picture(str(path), left, top, height=height)


def table_slide_rows(slide, headers, rows, left, top, col_widths, row_h=Inches(0.32), header_size=10, cell_size=9):
    y = top
    x = left
    for i, h in enumerate(headers):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[i], row_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = h
        tf.paragraphs[0].font.size = Pt(header_size)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += col_widths[i]
    for ri, row in enumerate(rows):
        y += row_h
        x = left
        fill = LIGHT if ri % 2 == 0 else WHITE
        for i, val in enumerate(row):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[i], row_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = RGBColor(0xD8, 0xDE, 0xE8)
            tf = shape.text_frame
            tf.paragraphs[0].text = str(val)
            tf.paragraphs[0].font.size = Pt(cell_size)
            tf.paragraphs[0].font.color.rgb = DARK
            tf.paragraphs[0].font.name = "Calibri"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER if i else PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            x += col_widths[i]


def overall(c: dict) -> float:
    s = c["themes_scores"]
    return sum(s.values()) / len(s)


def build() -> None:
    raw = json.loads(DATA_PATH.read_text(encoding="utf-8"))
    countries = sorted(raw["countries"], key=lambda c: -overall(c))
    trends = raw["global_trends"]
    maturity_counts = Counter(c["maturity_level"] for c in countries)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ----- 1 TITLE -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, NAVY)
    strip = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), Inches(13.33), Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = GOLD
    strip.line.fill.background()
    box(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5), "AI CLINIC", size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    box(
        sl,
        Inches(0.8),
        Inches(2.1),
        Inches(11.7),
        Inches(1.3),
        "AI for Healthcare Compliance and Regulations Across Countries",
        size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(3.5),
        Inches(11.7),
        Inches(0.6),
        "Interactive Dashboard and Comparative Regulatory Analysis",
        size=18,
        color=RGBColor(0xC8, 0xD4, 0xE8),
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(4.9),
        Inches(11.7),
        Inches(0.4),
        "Remi Uttejitha ALLAM  ·  Baptiste Langlois  ·  Tong Li  ·  Darryl Towa",
        size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(5.4),
        Inches(11.7),
        Inches(0.35),
        "Supervisor: Dr. Anuradha Kar  |  aivancity  |  Academic year 2025–2026",
        size=14,
        color=RGBColor(0xA8, 0xB8, 0xD0),
        align=PP_ALIGN.CENTER,
    )

    # ----- 2 AGENDA -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Agenda")
    bullets(
        sl,
        [
            "1. What is this project about?",
            "2. Problem statement and objectives",
            "3. Methodology pipeline (data → scores → dashboard)",
            "4. Compliance scoring and maturity classes",
            "5. Seven regulatory themes explained",
            "6. Key results — tables and figures",
            "7. Gaps, convergence, and recommendations",
            "8. Future impact — how this work helps improve practice",
            "9. Limitations and thank you",
        ],
        Inches(1.0),
        Inches(1.5),
        Inches(11),
        Inches(5.2),
        size=18,
        after=Pt(12),
    )

    # ----- 3 WHAT IS THE PROJECT -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "What is this project about?")
    box(
        sl,
        Inches(0.6),
        Inches(1.35),
        Inches(12.1),
        Inches(1.1),
        "AI Clinic compares how different countries regulate artificial intelligence in healthcare, and turns that comparison into an interactive dashboard and research report.",
        size=16,
        color=MED,
    )
    cards = [
        ("Literature synthesis", "Policy and peer-reviewed sources on SaMD, privacy, ethics, and lifecycle governance (2018–2026)."),
        ("20-country dataset", "Seven themes scored 1–10, maturity labels, device counts, challenges, and notable developments."),
        ("Interactive dashboard", "Map, theme analysis, peer comparison, trends, country profiles, and literature browser."),
        ("Research report", "Method, results figures/tables, country profiles, conclusions, and recommendations."),
    ]
    x = Inches(0.55)
    for title, desc in cards:
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.7), Inches(2.95), Inches(3.9))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = RGBColor(0xD0, 0xD6, 0xE0)
        box(sl, x + Inches(0.15), Inches(2.95), Inches(2.65), Inches(0.8), title, size=14, bold=True, color=NAVY)
        box(sl, x + Inches(0.15), Inches(3.75), Inches(2.65), Inches(2.5), desc, size=12, color=MED)
        x += Inches(3.15)

    # ----- 4 PROBLEM -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Problem: why does this matter?")
    bullets(
        sl,
        [
            "AI medical software (SaMD) is growing fast in imaging, clinical decision support, genomics, and drug discovery.",
            "Rules are not the same everywhere: privacy laws, approval pathways, validation standards, and liability differ by country.",
            "Public information is fragmented across statutes, agency guidance, ethics documents, and device lists.",
            "Developers, policymakers, clinicians, and researchers need one structured, comparable view.",
            "Our response: a coded comparative framework + live dashboard (decision support, not legal advice).",
        ],
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        size=17,
        after=Pt(14),
    )

    # ----- 5 OBJECTIVES -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Project objectives")
    bullets(
        sl,
        [
            "Synthesise literature and regulatory texts on AI healthcare compliance across 19 jurisdictions.",
            "Define seven thematic dimensions and score maturity on a transparent 1–10 scale.",
            "Build a curated, machine-readable dataset with regional metadata and global trends.",
            "Implement an interactive dashboard for exploration, comparison, and literature browsing.",
            "Document methodology, architecture, findings, and recommendations in the research report.",
        ],
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        size=17,
        after=Pt(14),
    )

    # ----- 6 PIPELINE -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Methodology pipeline: data collection → score calculation")
    steps = [
        ("1. Scope", "19 jurisdictions across seven regions: North America, Europe, Asia, Middle East, Africa, Oceania, South America."),
        ("2. Collect", "Primary: FDA, EMA, MHRA, NMPA, PMDA, WHO, IMDRF, GDPR/AI Act…\nSecondary: PubMed, Scopus, IEEE, Scholar, SSRN (2018–2026)."),
        ("3. Code", "Same schema for every country: privacy, AI rule, device framework, approval, validation, ethics, post-market, liability…"),
        ("4. Score", "Assign 1–10 per theme using 4 criteria; assign qualitative maturity label separately."),
        ("5. Store", "Version-controlled JSON: data/compliance_dataset.json."),
        ("6. Derive", "Dashboard computes overall mean, gaps (max−min), and use-case readiness."),
        ("7. Report", "Figures, tables, country profiles, conclusions."),
    ]
    y = Inches(1.35)
    for label, desc in steps:
        box(sl, Inches(0.55), y, Inches(1.7), Inches(0.7), label, size=13, bold=True, color=GOLD)
        box(sl, Inches(2.3), y, Inches(10.4), Inches(0.7), desc, size=13, color=DARK)
        y += Inches(0.75)

    # ----- 7 SCORING -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Compliance score — what it is and how it is calculated")
    bullets(
        sl,
        [
            "Each theme score is ordinal 1–10: 1 = little/no documented requirement; 10 = mature, enforceable practice.",
            "Evidence mapped with four criteria: (1) law/guidance exists, (2) AI/SaMD specificity, (3) enforcement/reporting, (4) IMDRF/WHO/GMLP alignment.",
            "Overall score = arithmetic mean of selected theme scores (rounded to 1 decimal).",
            "Formula: Overall = (1/k) × Σ sᵢ   where k = number of selected themes.",
            "Important: these are comparative research indicators — not official government ratings.",
            "No machine-learning model invents scores at runtime; scores are curated from documentary evidence.",
        ],
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.2),
        size=16,
        after=Pt(12),
    )

    # ----- 8 THEMES TABLE -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Seven dashboard variables (compliance themes)")
    table_slide_rows(
        sl,
        ["Theme", "What it measures"],
        [
            ["Data Privacy", "Health-data protection laws, cross-border transfers, DPIAs"],
            ["Clinical Validation", "Evidence standards, clinical evaluation, bias assessment"],
            ["Approval Process", "SaMD pathways, AI-specific guidance, regulatory sandboxes"],
            ["Transparency", "Explainability expectations, documentation, clinician intelligibility"],
            ["Ethics", "National AI ethics frameworks, fairness, human oversight"],
            ["Post-Market", "Adverse-event reporting, lifecycle monitoring, update governance"],
            ["Liability", "Product liability, malpractice interfaces, AI-specific proposals"],
        ],
        Inches(0.7),
        Inches(1.4),
        [Inches(3.2), Inches(8.8)],
        row_h=Inches(0.68),
        header_size=13,
        cell_size=13,
    )

    # ----- 9 MATURITY -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Maturity classes — definitions (not a formula on overall score)")
    table_slide_rows(
        sl,
        ["Class", "Definition", "Examples", "n"],
        [
            ["Early", "Limited AI-specific healthcare rules; few approvals", "Nigeria, Kenya", str(maturity_counts.get("Early", 0))],
            ["Emerging", "Privacy/ethics exist; SaMD pathway still forming", "Saudi Arabia, South Africa", str(maturity_counts.get("Emerging", 0))],
            ["Developing", "Device framework + emerging AI guidance", "India, Brazil", str(maturity_counts.get("Developing", 0))],
            ["Moderate", "Working SaMD registration & post-market system", "Australia", str(maturity_counts.get("Moderate", 0))],
            ["Advanced", "Dedicated AI/ML SaMD or high-risk AI law + lifecycle", "US, EU, Japan, Singapore…", str(maturity_counts.get("Advanced", 0))],
        ],
        Inches(0.45),
        Inches(1.4),
        [Inches(1.7), Inches(6.0), Inches(3.8), Inches(0.9)],
        row_h=Inches(0.72),
        header_size=12,
        cell_size=12,
    )
    box(
        sl,
        Inches(0.55),
        Inches(5.5),
        Inches(12.2),
        Inches(1.2),
        "Maturity is assigned from documentary evidence (AI rules, device pathway, approval activity, privacy/ethics enforcement). Device risk class (FDA I–III) describes the product; maturity describes the country.",
        size=13,
        color=MED,
    )

    # ----- 10 FIGURE overall ranking -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Figure: overall country ranking")
    add_picture_safe(sl, FIG_DIR / "overall_ranking.png", Inches(1.8), Inches(1.3), width=Inches(9.7))
    box(
        sl,
        Inches(0.55),
        Inches(6.55),
        Inches(12.2),
        Inches(0.5),
        "Mean of seven theme scores. EU / United States lead; Kenya and Nigeria trail — capacity gap at a glance.",
        size=12,
        color=MED,
    )

    # ----- 11 FIGURE themes -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Figures: theme averages and inter-country gaps")
    add_picture_safe(sl, FIG_DIR / "theme_averages.png", Inches(0.35), Inches(1.35), width=Inches(6.2))
    add_picture_safe(sl, FIG_DIR / "theme_gaps.png", Inches(6.7), Inches(1.35), width=Inches(6.2))
    box(
        sl,
        Inches(0.55),
        Inches(6.4),
        Inches(12.2),
        Inches(0.55),
        "Left: worldwide average maturity by theme (approval & privacy lead; post-market & liability lag). Right: Gap = max − min score — large gaps mean strong divergence.",
        size=12,
        color=MED,
    )

    # ----- 12 FIGURE maturity + regional -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Figures: maturity mix and regional means")
    add_picture_safe(sl, FIG_DIR / "maturity_mix.png", Inches(0.4), Inches(1.35), width=Inches(5.8))
    add_picture_safe(sl, FIG_DIR / "regional_overall.png", Inches(6.5), Inches(1.5), width=Inches(6.3))
    box(
        sl,
        Inches(0.55),
        Inches(6.45),
        Inches(12.2),
        Inches(0.5),
        "Advanced slice is large because the sample includes major regulators. Europe leads regional means; Africa is lowest.",
        size=12,
        color=MED,
    )

    # ----- 13 FIGURE peer + radar -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Figures: peer comparison and top-5 radar profiles")
    add_picture_safe(sl, FIG_DIR / "peer_comparison.png", Inches(0.3), Inches(1.35), width=Inches(7.0))
    add_picture_safe(sl, FIG_DIR / "radar_top5.png", Inches(7.4), Inches(1.4), width=Inches(5.5))
    box(
        sl,
        Inches(0.55),
        Inches(6.45),
        Inches(12.2),
        Inches(0.5),
        "US/EU lead most themes; India and South Africa lag on validation, post-market, and liability. Radar shape shows theme strengths differ even among top scorers.",
        size=12,
        color=MED,
    )

    # ----- 14 SCORES TABLE -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Table: country compliance scores (top 12 + bottom 3)")
    headers = ["Country", "Overall", "Priv", "Clin", "Appr", "Trans", "Eth", "Post", "Liab", "Maturity"]
    widths = [Inches(2.2), Inches(0.95)] + [Inches(0.85)] * 7 + [Inches(1.5)]
    show = countries[:12] + countries[-3:]
    rows = []
    for c in show:
        s = c["themes_scores"]
        rows.append(
            [
                c["country"][:18],
                f"{overall(c):.1f}",
                s["data_privacy"],
                s["clinical_validation"],
                s["approval_process"],
                s["transparency"],
                s["ethics"],
                s["post_market"],
                s["liability"],
                c["maturity_level"][:10],
            ]
        )
    table_slide_rows(sl, headers, rows, Inches(0.35), Inches(1.25), widths, row_h=Inches(0.33), header_size=10, cell_size=10)

    # ----- 15 KEY RESULTS -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Headline results")
    bullets(
        sl,
        [
            f"Sample: {len(countries)} jurisdictions · 7 themes · maturity distribution includes {maturity_counts.get('Advanced', 0)} Advanced and {maturity_counts.get('Early', 0)} Early.",
            f"Top overall: {countries[0]['country']} ({overall(countries[0]):.1f}), {countries[1]['country']} ({overall(countries[1]):.1f}), {countries[2]['country']} ({overall(countries[2]):.1f}).",
            "Europe leads privacy/ethics/transparency; United States leads clinical validation and approval volume.",
            "Asia is heterogeneous: Japan/Singapore/China Advanced; India/Thailand Developing.",
            "Africa: South Africa Emerging; Kenya & Nigeria Early — privacy statutes alone do not raise overall AI healthcare compliance.",
            "Weakest global pillars: post-market surveillance and liability (lifecycle & accountability).",
        ],
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.3),
        size=16,
        after=Pt(12),
    )

    # ----- 16 GAPS & CONVERGENCE -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "What we mean by gaps and convergence trends")
    left = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.2))
    left.fill.solid()
    left.fill.fore_color.rgb = LIGHT
    left.line.fill.background()
    box(sl, Inches(0.75), Inches(1.6), Inches(5.4), Inches(0.45), "Gaps", size=20, bold=True, color=NAVY)
    bullets(
        sl,
        [
            "Gap = max(country score) − min(country score) on one theme.",
            "Large gap = divergence zone.",
            "Useful for multi-country evidence strategy.",
            "Widest gaps often: clinical validation & transparency.",
        ],
        Inches(0.75),
        Inches(2.3),
        Inches(5.4),
        Inches(3.8),
        size=14,
        after=Pt(10),
    )
    right = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.4), Inches(5.9), Inches(5.2))
    right.fill.solid()
    right.fill.fore_color.rgb = LIGHT
    right.line.fill.background()
    box(sl, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.45), "Convergence", size=20, bold=True, color=NAVY)
    bullets(
        sl,
        [
            "Not identical laws — shared principles & tools.",
            "Risk-based SaMD classification (IMDRF-style).",
            "GDPR-like health-data protection.",
            "GMLP / lifecycle / PCCP thinking.",
            "Shift from voluntary ethics toward binding rules (e.g. EU AI Act).",
        ],
        Inches(7.1),
        Inches(2.3),
        Inches(5.4),
        Inches(3.8),
        size=14,
        after=Pt(10),
    )

    # ----- 17 TRENDS TABLE (first 8) -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Table: selected global trends")
    trows = []
    for t in trends[:8]:
        trows.append([t["trend"][:42], t["adoption_level"], str(t["year_emerged"]), t["description"][:70] + "…"])
    table_slide_rows(
        sl,
        ["Trend", "Adoption", "Since", "Description (short)"],
        trows,
        Inches(0.4),
        Inches(1.35),
        [Inches(3.6), Inches(1.2), Inches(1.0), Inches(6.5)],
        row_h=Inches(0.55),
        header_size=11,
        cell_size=10,
    )

    # ----- 18 DASHBOARD -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Interactive dashboard artefact")
    bullets(
        sl,
        [
            "Public web app (React + TypeScript + Recharts + maps), hosted on GitHub Pages.",
            "Tabs: Overview · Theme Analysis · Comparison · Trends & Use Cases · Country Detail · Literature.",
            "Filters by region, maturity, and themes recalculate overall scores and gaps instantly.",
            "Live URL: https://remi7025.github.io/AI-Healthcare-Compliance-Dashboard/",
            "Canonical data: data/compliance_dataset.json (same source as this presentation’s tables).",
        ],
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        size=16,
        after=Pt(14),
    )

    # ----- 19 FUTURE / HOW IT HELPS -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Future impact — how this work helps improve")
    bullets(
        sl,
        [
            "For developers: choose first markets by theme strength (e.g. validation-heavy imaging → US/EU).",
            "For policymakers: see where gaps are largest (validation, post-market, liability) and prioritise capacity building.",
            "For researchers & students: reproducible open dataset + dashboard for teaching comparative AI governance.",
            "For hospitals/clinics: link use-case readiness (radiology, genomics…) to local regulatory maturity.",
            "Future upgrades: live regulatory feeds, Delphi expert scoring, more Latin America & African regulators,",
            "and a dedicated evaluation track for foundation-model medical products.",
            "Net effect: faster, safer, more transparent market access decisions and clearer reform priorities.",
        ],
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.3),
        size=15,
        after=Pt(11),
    )

    # ----- 20 RECOMMENDATIONS -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Recommendations (synthesised from WHO, IMDRF, GMLP, EU AI Act, NIST, literature)")
    bullets(
        sl,
        [
            "1. Expand IMDRF/WHO coordination on AI-specific clinical-evidence methods.",
            "2. Build regulatory capacity in developing countries (technical assistance, mutual-recognition pilots).",
            "3. Operationalise ethics through measurable fairness, monitoring, and documentation standards.",
            "4. Update liability frameworks before harm cases proliferate.",
            "5. Maintain public dashboards and open datasets for transparent comparative research.",
            "6. Treat foundation-model medical products as a dedicated evaluation track.",
            "Note: recommendations are a synthesis of the evidence base — not a new Delphi survey.",
        ],
        Inches(0.7),
        Inches(1.4),
        Inches(12),
        Inches(5.3),
        size=15,
        after=Pt(11),
    )

    # ----- 21 LIMITATIONS -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, "Limitations")
    bullets(
        sl,
        [
            "Scores are expert documentary syntheses and will drift as laws change.",
            "Device-approval counts are incomplete where agencies do not publish AI-specific tallies.",
            "Sample over-represents Advanced jurisdictions relative to the true world map.",
            "Oceania and South America regional means rest on single countries (Australia, Brazil).",
            "This is academic decision support — not legal advice for market authorisation.",
        ],
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        size=16,
        after=Pt(14),
    )

    # ----- 22 THANK YOU -----
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, NAVY)
    gold_bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.85), Inches(4.3), Inches(0.08))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()
    box(sl, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.9), "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(
        sl,
        Inches(0.8),
        Inches(3.2),
        Inches(11.7),
        Inches(0.5),
        "Questions & Discussion",
        size=22,
        color=GOLD,
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(4.2),
        Inches(11.7),
        Inches(0.4),
        "Remi Uttejitha ALLAM  ·  Baptiste Langlois  ·  Tong Li  ·  Darryl Towa",
        size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(4.75),
        Inches(11.7),
        Inches(0.35),
        "Supervisor: Dr. Anuradha Kar  |  aivancity School for Technology, Business & Society",
        size=13,
        color=RGBColor(0xA8, 0xB8, 0xD0),
        align=PP_ALIGN.CENTER,
    )
    box(
        sl,
        Inches(0.8),
        Inches(5.4),
        Inches(11.7),
        Inches(0.35),
        "Dashboard: remi7025.github.io/AI-Healthcare-Compliance-Dashboard",
        size=12,
        color=RGBColor(0x90, 0xA4, 0xC0),
        align=PP_ALIGN.CENTER,
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
